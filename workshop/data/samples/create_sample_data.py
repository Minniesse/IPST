"""
สร้างไฟล์ตัวอย่างภาษาไทยสำหรับ Workshop
Create sample Thai PDF and PPTX files for the workshop.

Usage:
    uv run python data/samples/create_sample_data.py
"""

from pathlib import Path


THAI_CONTENT = {
    "title": "ปัญญาประดิษฐ์และการเรียนรู้ของเครื่อง",
    "sections": [
        {
            "heading": "บทที่ 1: ปัญญาประดิษฐ์คืออะไร",
            "text": (
                "ปัญญาประดิษฐ์ (Artificial Intelligence หรือ AI) คือสาขาวิชาทางวิทยาการคอมพิวเตอร์ "
                "ที่มุ่งเน้นการสร้างระบบที่สามารถทำงานที่ต้องอาศัยความฉลาดของมนุษย์ เช่น "
                "การเรียนรู้ การตัดสินใจ การรู้จำภาพ การเข้าใจภาษา และการแก้ปัญหา\n\n"
                "AI สามารถแบ่งออกเป็นหลายประเภท ได้แก่:\n"
                "1. AI แบบแคบ (Narrow AI) - ออกแบบมาเพื่อทำงานเฉพาะอย่าง เช่น การแปลภาษา\n"
                "2. AI แบบทั่วไป (General AI) - สามารถทำงานได้หลากหลายเหมือนมนุษย์\n"
                "3. AI ขั้นสูง (Super AI) - มีความสามารถเหนือกว่ามนุษย์ในทุกด้าน"
            ),
        },
        {
            "heading": "บทที่ 2: การเรียนรู้ของเครื่อง",
            "text": (
                "การเรียนรู้ของเครื่อง (Machine Learning) เป็นสาขาย่อยของปัญญาประดิษฐ์ "
                "ที่ทำให้คอมพิวเตอร์สามารถเรียนรู้จากข้อมูลได้โดยไม่ต้องเขียนโปรแกรมอย่างชัดเจน\n\n"
                "ประเภทหลักของการเรียนรู้ของเครื่อง:\n"
                "- การเรียนรู้แบบมีผู้สอน (Supervised Learning) - เรียนรู้จากข้อมูลที่มีคำตอบ\n"
                "- การเรียนรู้แบบไม่มีผู้สอน (Unsupervised Learning) - หาโครงสร้างจากข้อมูลเอง\n"
                "- การเรียนรู้แบบเสริมกำลัง (Reinforcement Learning) - เรียนรู้จากรางวัลและบทลงโทษ\n\n"
                "ตัวอย่างการใช้งาน Machine Learning ในชีวิตประจำวัน ได้แก่ "
                "ระบบแนะนำสินค้า ระบบกรองอีเมลขยะ การรู้จำใบหน้า "
                "ผู้ช่วยเสียง และรถยนต์ขับเคลื่อนอัตโนมัติ"
            ),
        },
        {
            "heading": "บทที่ 3: การเรียนรู้เชิงลึก",
            "text": (
                "การเรียนรู้เชิงลึก (Deep Learning) เป็นเทคนิคขั้นสูงของ Machine Learning "
                "ที่ใช้โครงข่ายประสาทเทียม (Neural Networks) หลายชั้นในการเรียนรู้\n\n"
                "โครงข่ายประสาทเทียมประกอบด้วย:\n"
                "- ชั้นรับข้อมูล (Input Layer) - รับข้อมูลดิบเข้าสู่ระบบ\n"
                "- ชั้นซ่อน (Hidden Layers) - ประมวลผลและสกัดคุณลักษณะ\n"
                "- ชั้นผลลัพธ์ (Output Layer) - ให้คำตอบหรือการทำนาย\n\n"
                "สถาปัตยกรรมที่สำคัญ:\n"
                "1. CNN (Convolutional Neural Network) - เหมาะสำหรับงานรูปภาพ\n"
                "2. RNN (Recurrent Neural Network) - เหมาะสำหรับข้อมูลลำดับ\n"
                "3. Transformer - สถาปัตยกรรมที่ใช้ใน GPT, BERT และโมเดลภาษาขนาดใหญ่"
            ),
        },
        {
            "heading": "บทที่ 4: โมเดลภาษาขนาดใหญ่",
            "text": (
                "โมเดลภาษาขนาดใหญ่ (Large Language Models หรือ LLM) เป็นโมเดล AI "
                "ที่ได้รับการฝึกฝนด้วยข้อมูลข้อความจำนวนมหาศาล "
                "ทำให้สามารถเข้าใจและสร้างข้อความภาษามนุษย์ได้อย่างเป็นธรรมชาติ\n\n"
                "การเตรียมข้อมูลสำหรับสอนโมเดลมีขั้นตอนหลักดังนี้:\n"
                "1. การรวบรวมข้อมูล (Data Collection) - เก็บข้อมูลจากหลายแหล่ง\n"
                "2. การทำความสะอาดข้อมูล (Data Cleaning) - ลบข้อมูลที่ไม่ต้องการ\n"
                "3. การจัดโครงสร้างข้อมูล (Data Structuring) - จัดรูปแบบให้เหมาะกับการเรียนรู้\n"
                "4. การตรวจสอบคุณภาพ (Quality Assurance) - ตรวจสอบความถูกต้องของข้อมูล\n\n"
                "ตัวอย่างโมเดลภาษาขนาดใหญ่ที่สำคัญ:\n"
                "- GPT-4 จาก OpenAI\n"
                "- Claude จาก Anthropic\n"
                "- Gemini จาก Google\n"
                "- Llama จาก Meta\n"
                "- Qwen จาก Alibaba"
            ),
        },
    ],
    "table_data": {
        "title": "เปรียบเทียบประเภทของ Machine Learning",
        "headers": ["ประเภท", "ข้อมูลที่ใช้", "ตัวอย่างการใช้งาน", "ข้อดี"],
        "rows": [
            ["Supervised Learning", "มีคำตอบ (Labels)", "การจำแนกรูปภาพ", "แม่นยำสูง"],
            ["Unsupervised Learning", "ไม่มีคำตอบ", "การจัดกลุ่มลูกค้า", "ค้นพบโครงสร้างใหม่"],
            ["Reinforcement Learning", "รางวัล/บทลงโทษ", "เกม, หุ่นยนต์", "ปรับตัวได้ดี"],
        ],
    },
}


def create_sample_pdf(output_path: Path) -> None:
    """Create a sample Thai PDF using Docling-compatible markdown-to-PDF approach."""
    # Write as a well-structured Markdown file that Docling can also process
    md_content = f"# {THAI_CONTENT['title']}\n\n"

    for section in THAI_CONTENT["sections"]:
        md_content += f"## {section['heading']}\n\n"
        md_content += f"{section['text']}\n\n"

    # Add table
    table = THAI_CONTENT["table_data"]
    md_content += f"## {table['title']}\n\n"
    md_content += "| " + " | ".join(table["headers"]) + " |\n"
    md_content += "| " + " | ".join(["---"] * len(table["headers"])) + " |\n"
    for row in table["rows"]:
        md_content += "| " + " | ".join(row) + " |\n"
    md_content += "\n"

    # Save as Markdown (Docling can process .md files too)
    md_path = output_path / "thai_sample.md"
    md_path.write_text(md_content, encoding="utf-8")
    print(f"Created: {md_path}")

    # Also try to create PDF using a simple approach
    try:
        from docling_core.types.doc import DoclingDocument, TableData, TableCell

        doc = DoclingDocument(name="thai_sample")

        for section in THAI_CONTENT["sections"]:
            doc.add_heading(text=section["heading"])
            for paragraph in section["text"].split("\n\n"):
                if paragraph.strip():
                    doc.add_text(text=paragraph.strip(), label="paragraph")

        # Save as JSON that we can use
        json_path = output_path / "thai_sample_docling.json"
        import json
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(doc.export_to_dict(), f, ensure_ascii=False, indent=2)
        print(f"Created: {json_path}")
    except ImportError:
        print("docling_core not available, skipping DoclingDocument creation")


def create_sample_pptx(output_path: Path) -> None:
    """Create a sample Thai PPTX file."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = THAI_CONTENT["title"]
        slide.placeholders[1].text = "Workshop: การเตรียมข้อมูลสำหรับสอนโมเดล"

        # Content slides
        for section in THAI_CONTENT["sections"]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = section["heading"]
            body = slide.placeholders[1]
            body.text = section["text"]

        # Table slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
        slide.shapes.title.text = THAI_CONTENT["table_data"]["title"]

        table_data = THAI_CONTENT["table_data"]
        rows = len(table_data["rows"]) + 1  # +1 for header
        cols = len(table_data["headers"])

        table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(3))
        table = table_shape.table

        # Header row
        for j, header in enumerate(table_data["headers"]):
            table.cell(0, j).text = header

        # Data rows
        for i, row in enumerate(table_data["rows"]):
            for j, cell_text in enumerate(row):
                table.cell(i + 1, j).text = cell_text

        pptx_path = output_path / "thai_sample.pptx"
        prs.save(str(pptx_path))
        print(f"Created: {pptx_path}")

    except ImportError:
        print("python-pptx not available, skipping PPTX creation")


def create_sample_pdf_from_md(output_path: Path) -> None:
    """Create a simple PDF from the markdown content using basic fpdf2."""
    try:
        from fpdf import FPDF

        pdf = FPDF()
        pdf.add_page()

        # Try to add Thai font
        font_paths = [
            "/usr/share/fonts/truetype/noto/NotoSansThai-Regular.ttf",
            "/usr/share/fonts/truetype/thai/TlwgTypo.ttf",
            "/usr/share/fonts/truetype/tlwg/TlwgTypo.ttf",
        ]
        thai_font_added = False
        for font_path in font_paths:
            if Path(font_path).exists():
                pdf.add_font("Thai", "", font_path, uni=True)
                pdf.set_font("Thai", size=12)
                thai_font_added = True
                break

        if not thai_font_added:
            pdf.set_font("Helvetica", size=12)
            print("Warning: No Thai font found, PDF may not display Thai correctly")

        pdf.cell(0, 10, THAI_CONTENT["title"], new_x="LMARGIN", new_y="NEXT")
        pdf.ln(5)

        for section in THAI_CONTENT["sections"]:
            pdf.set_font_size(14)
            pdf.cell(0, 10, section["heading"], new_x="LMARGIN", new_y="NEXT")
            pdf.set_font_size(11)
            pdf.multi_cell(0, 6, section["text"])
            pdf.ln(5)

        pdf_path = output_path / "thai_sample.pdf"
        pdf.output(str(pdf_path))
        print(f"Created: {pdf_path}")

    except ImportError:
        print("fpdf2 not available, skipping PDF creation")
        print("Install with: uv add fpdf2")


if __name__ == "__main__":
    output_path = Path(__file__).parent
    output_path.mkdir(parents=True, exist_ok=True)

    print("Creating sample Thai documents for the workshop...\n")

    # Always create markdown version (Docling can process it)
    create_sample_pdf(output_path)

    # Create PPTX
    create_sample_pptx(output_path)

    # Try to create PDF
    create_sample_pdf_from_md(output_path)

    print("\nDone! Sample files are in:", output_path)
